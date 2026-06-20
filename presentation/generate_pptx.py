"""Run this script once to produce teyavet_presentation.pptx in the same directory."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── palette ──────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x11, 0x17)
SURFACE  = RGBColor(0x16, 0x1b, 0x22)
ACCENT   = RGBColor(0x58, 0xa6, 0xff)  # blue
GREEN    = RGBColor(0x3f, 0xb9, 0x50)
RED      = RGBColor(0xf7, 0x81, 0x66)
PURPLE   = RGBColor(0xd2, 0xa8, 0xff)
TEXT     = RGBColor(0xe6, 0xed, 0xf3)
MUTED    = RGBColor(0x8b, 0x94, 0x9e)
BORDER   = RGBColor(0x30, 0x36, 0x3d)

OUT_FILE = os.path.join(os.path.dirname(__file__), "teyavet_presentation.pptx")

# ── helpers ───────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, prs, color=BG):
    from pptx.util import Emu
    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def add_rect(slide, l, t, w, h, fill=SURFACE, line=BORDER, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = line_w
    return shape

def add_text(slide, text, l, t, w, h,
             size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(tf, text, size=12, bold=False, color=TEXT, bullet=True, indent=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    if bullet:
        p.text = ("  " * indent) + "› " + text
    else:
        p.text = ("  " * indent) + text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return p

def add_divider(slide, t, w=Inches(12.5), l=Inches(0.42)):
    line = slide.shapes.add_shape(1, l, t, w, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def slide_counter(slide, prs, n, total):
    add_text(slide, f"{n} / {total}",
             prs.slide_width - Inches(1.2),
             prs.slide_height - Inches(0.45),
             Inches(1.1), Inches(0.35),
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)

def slide_label(slide, label):
    add_text(slide, label.upper(),
             Inches(0.42), Inches(0.22),
             Inches(8), Inches(0.28),
             size=9, bold=True, color=MUTED)

# ── SLIDE 1 ───────────────────────────────────────────────

def make_slide1(prs):
    slide = blank_slide(prs)
    fill_bg(slide, prs)

    slide_label(slide, "DevSecOps — Final Project")

    # Title
    add_text(slide, "TeyaVet — Vet Clinic Management System",
             Inches(0.42), Inches(0.5), Inches(12), Inches(0.55),
             size=28, bold=True, color=ACCENT)

    add_divider(slide, Inches(1.18))

    # ── left: name card ──
    card_l = add_rect(slide, Inches(0.42), Inches(1.35), Inches(2.8), Inches(2.5))
    tb = slide.shapes.add_textbox(Inches(0.52), Inches(1.45), Inches(2.6), Inches(2.3))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "AM"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = ACCENT

    for line, sz, col in [
        ("Assaf Mashal",              13, TEXT),
        ("Full-Stack & Cloud Engineer", 10, MUTED),
        ("DevSecOps Course — Final Project", 9, ACCENT),
    ]:
        add_para(tf, line, size=sz, color=col, bullet=False)
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER

    # ── right top: What is TeyaVet ──
    add_rect(slide, Inches(3.52), Inches(1.35), Inches(4.5), Inches(1.5))
    add_text(slide, "WHAT IS TEYAVET?",
             Inches(3.65), Inches(1.42), Inches(4), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb2 = slide.shapes.add_textbox(Inches(3.65), Inches(1.72), Inches(4.2), Inches(1.1))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for item in [
        "Vet clinic staff portal",
        "Manage pet owners, patients & appointments",
        "Scheduling with conflict detection",
        "Medical records & measurements",
    ]:
        add_para(tf2, item, size=11, color=TEXT)

    # ── right bottom: Tech stack ──
    add_rect(slide, Inches(8.22), Inches(1.35), Inches(4.5), Inches(1.5))
    add_text(slide, "TECH STACK",
             Inches(8.35), Inches(1.42), Inches(4), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb3 = slide.shapes.add_textbox(Inches(8.35), Inches(1.72), Inches(4.2), Inches(1.1))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    for item in [
        "Backend — Python 3.13 · Flask · JWT",
        "Frontend — React 18 · Vite · plain CSS",
        "Database — MySQL 8 (InnoDB)",
        "Infra — Docker → Kubernetes on AWS",
    ]:
        add_para(tf3, item, size=11, color=TEXT)

    # ── bottom: architecture card ──
    add_rect(slide, Inches(0.42), Inches(4.05), Inches(12.5), Inches(1.9))
    add_text(slide, "ARCHITECTURE — DEV VS PRODUCTION",
             Inches(0.55), Inches(4.12), Inches(8), Inches(0.28),
             size=8, bold=True, color=MUTED)

    # Dev row
    add_text(slide, "LOCAL / DEV",
             Inches(0.55), Inches(4.45), Inches(2), Inches(0.25),
             size=8, bold=True, color=MUTED)
    add_text(slide, "docker compose up  →  MySQL 8  +  Flask API :5000  +  React :5173",
             Inches(0.55), Inches(4.72), Inches(5.8), Inches(0.3),
             size=11, color=ACCENT)

    # Prod row
    add_text(slide, "PRODUCTION (AWS)",
             Inches(6.8), Inches(4.45), Inches(2.5), Inches(0.25),
             size=8, bold=True, color=MUTED)
    add_text(slide, "ALB  →  EKS pods (backend + frontend)  +  RDS MySQL (private subnet)  ·  ECR images",
             Inches(6.8), Inches(4.72), Inches(6.0), Inches(0.3),
             size=11, color=GREEN)

    slide_counter(slide, prs, 1, 3)
    return slide


# ── SLIDE 2 ───────────────────────────────────────────────

def make_slide2(prs):
    slide = blank_slide(prs)
    fill_bg(slide, prs)

    slide_label(slide, "Infrastructure & Hosting")
    add_text(slide, "What Runs Where — AWS Architecture",
             Inches(0.42), Inches(0.5), Inches(12), Inches(0.55),
             size=28, bold=True, color=ACCENT)
    add_divider(slide, Inches(1.18))

    # ── left column ──
    add_rect(slide, Inches(0.42), Inches(1.35), Inches(6.1), Inches(2.65))
    add_text(slide, "AWS RESOURCES (TERRAFORM IaC)",
             Inches(0.55), Inches(1.42), Inches(5.8), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.72), Inches(5.8), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    items = [
        ("VPC", " — 2 public + 2 private subnets, multi-AZ"),
        ("EKS 1.31", " — managed node group (t3.medium, private subnets)"),
        ("RDS MySQL 8", " — private subnet only, no public access"),
        ("ECR", " — 2 repos (backend + frontend), image scanning on push"),
        ("ALB", " — provisioned by AWS Load Balancer Controller via Helm"),
    ]
    for key, val in items:
        p = tf.add_paragraph()
        p.font.size = Pt(11)
        r1 = p.add_run(); r1.text = "› " + key; r1.font.bold = True; r1.font.color.rgb = ACCENT
        r2 = p.add_run(); r2.text = val;          r2.font.color.rgb = TEXT

    add_text(slide, "Terraform  ·  S3 remote state  ·  DynamoDB lock",
             Inches(0.55), Inches(4.07), Inches(5.8), Inches(0.28),
             size=9, bold=True, color=GREEN)

    # ── left bottom ──
    add_rect(slide, Inches(0.42), Inches(4.22), Inches(6.1), Inches(1.85))
    add_text(slide, "WHERE ARE DOCKER IMAGES BUILT?",
             Inches(0.55), Inches(4.29), Inches(5.8), Inches(0.28),
             size=8, bold=True, color=MUTED)
    add_text(slide, "push to master  →  GitHub Actions runner  →  AWS ECR",
             Inches(0.55), Inches(4.58), Inches(5.8), Inches(0.3),
             size=12, bold=True, color=ACCENT)
    tb2 = slide.shapes.add_textbox(Inches(0.55), Inches(4.92), Inches(5.8), Inches(1.1))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for item in [
        "3 images: backend · frontend (multi-stage Node→Nginx) · db",
        "Backend: non-root appuser — least privilege",
        "ECR lifecycle: keep last 10 tags, expire untagged after 1 day",
    ]:
        add_para(tf2, item, size=10.5, color=TEXT)

    # ── right column ──
    add_rect(slide, Inches(6.72), Inches(1.35), Inches(6.2), Inches(2.65))
    add_text(slide, "RESPONSIBLE HOSTING — SECURITY CONTROLS",
             Inches(6.85), Inches(1.42), Inches(5.9), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb3 = slide.shapes.add_textbox(Inches(6.85), Inches(1.72), Inches(5.9), Inches(2.2))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    for item in [
        ("No static AWS keys", " — GitHub OIDC → STS short-lived tokens"),
        ("RDS in private subnet", " — unreachable from internet"),
        ("ECR image scanning", " on every push"),
        ("Non-root containers", " — Flask runs as appuser"),
        ("Secrets in GitHub Secrets", " — injected at runtime, never in code"),
        ("Security groups", ": RDS only accepts MySQL from EKS nodes"),
    ]:
        p = tf3.add_paragraph()
        p.font.size = Pt(11)
        r1 = p.add_run(); r1.text = "› " + item[0]; r1.font.bold = True; r1.font.color.rgb = GREEN
        r2 = p.add_run(); r2.text = item[1];          r2.font.color.rgb = TEXT

    add_rect(slide, Inches(6.72), Inches(4.22), Inches(6.2), Inches(1.85))
    add_text(slide, "IaC WORKFLOW",
             Inches(6.85), Inches(4.29), Inches(5.9), Inches(0.28),
             size=8, bold=True, color=MUTED)
    add_text(slide, "Bootstrap → S3+DynamoDB state → terraform-apply.yml → VPC / EKS / RDS / ECR",
             Inches(6.85), Inches(4.58), Inches(5.9), Inches(0.35),
             size=10.5, color=ACCENT)
    tb4 = slide.shapes.add_textbox(Inches(6.85), Inches(4.95), Inches(5.9), Inches(1.0))
    tf4 = tb4.text_frame; tf4.word_wrap = True
    for item in [
        "Manual dispatch — intentional gate before infra changes",
        "prevent_destroy on GitHub OIDC provider",
        "destroy.sh script for full teardown with safeguards",
    ]:
        add_para(tf4, item, size=10.5, color=TEXT)

    slide_counter(slide, prs, 2, 4)
    return slide


# ── SLIDE 3 ───────────────────────────────────────────────

def make_slide3(prs):
    slide = blank_slide(prs)
    fill_bg(slide, prs)

    slide_label(slide, "CI/CD & Security")
    add_text(slide, "DevSecOps Workflow — Git Branching & Pipelines",
             Inches(0.42), Inches(0.5), Inches(12), Inches(0.55),
             size=28, bold=True, color=ACCENT)
    add_divider(slide, Inches(1.18))

    # ── card 1: Git branching ──
    add_rect(slide, Inches(0.42), Inches(1.35), Inches(3.9), Inches(3.5))
    add_text(slide, "GIT BRANCH STRATEGY",
             Inches(0.55), Inches(1.42), Inches(3.7), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.72), Inches(3.7), Inches(3.0))
    tf = tb.text_frame; tf.word_wrap = True
    add_para(tf, "master = production (trunk-based)", size=11, color=TEXT, bullet=True)
    add_para(tf, "Feature branches by prefix:", size=11, color=TEXT, bullet=True)

    for name, col in [("feat/", ACCENT), ("ci/", GREEN), ("fix/", RED), ("infra/", PURPLE), ("security/", RED)]:
        p = tf.add_paragraph()
        p.font.size = Pt(11)
        r = p.add_run(); r.text = "    • " + name; r.font.bold = True; r.font.color.rgb = col

    add_para(tf, "All changes: PR → checks → merge", size=11, color=TEXT, bullet=True)
    add_para(tf, "No direct push to master", size=11, color=TEXT, bullet=True)

    # ── card 2: PR Gate ──
    add_rect(slide, Inches(4.52), Inches(1.35), Inches(4.1), Inches(3.5))
    add_text(slide, "PR GATE — pr-checks.yml",
             Inches(4.65), Inches(1.42), Inches(3.9), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb2 = slide.shapes.add_textbox(Inches(4.65), Inches(1.72), Inches(3.9), Inches(3.0))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    items2 = [
        (GREEN,  "Lint", " — ESLint (frontend) · ruff (backend)"),
        (GREEN,  "SCA",  " — npm audit --audit-level=high"),
        (GREEN,  "SCA",  " — pip-audit (Python deps)"),
        (GREEN,  "Unit tests", " — pytest backend suite"),
        (TEXT,   "Build must succeed", " (Vite build)"),
        (MUTED,  "Triggers", ": PR open / sync / reopen to master"),
    ]
    for col, key, val in items2:
        p = tf2.add_paragraph(); p.font.size = Pt(11)
        r1 = p.add_run(); r1.text = "› " + key; r1.font.bold = True; r1.font.color.rgb = col
        r2 = p.add_run(); r2.text = val; r2.font.color.rgb = TEXT

    # ── card 3: CI/CD pipeline (cd.yml only, ends at ECR) ──
    add_rect(slide, Inches(8.82), Inches(1.35), Inches(4.1), Inches(3.0))
    add_text(slide, "CI/CD — cd.yml (push to master)",
             Inches(8.95), Inches(1.42), Inches(3.9), Inches(0.28),
             size=8, bold=True, color=MUTED)
    steps = [
        (PURPLE, "push to master"),
        (ACCENT, "Build 3 Docker images"),
        (ACCENT, "Containers → integration tests"),
        (RED,    "Trivy scan — block on CRITICAL"),
        (GREEN,  "Push to AWS ECR (OIDC)"),
    ]
    y = Inches(1.72)
    for i, (col, text) in enumerate(steps):
        add_text(slide, text, Inches(8.95), y, Inches(3.9), Inches(0.30), size=10.5, color=col)
        y += Inches(0.33)
        if i < len(steps) - 1:
            add_text(slide, "↓", Inches(9.15), y - Inches(0.14),
                     Inches(0.3), Inches(0.20), size=9, color=MUTED)

    # ── bottom: EKS deploy + OIDC ──
    add_rect(slide, Inches(0.42), Inches(4.55), Inches(12.5), Inches(1.65))
    add_text(slide, "EKS DEPLOY — deploy-eks.yml   ·   manual workflow_dispatch   ·   environment: production (reviewer approval)",
             Inches(0.55), Inches(4.62), Inches(12.0), Inches(0.28),
             size=8, bold=True, color=MUTED)

    flow_items = [
        (ACCENT, "image_tag input"),
        (ACCENT, "OIDC → kubeconfig"),
        (PURPLE, "Helm: AWS LBC"),
        (RED,    "K8s Secret ← GitHub Secrets"),
        (ACCENT, "DB migration job"),
        (ACCENT, "Apply Services + Ingress"),
        (GREEN,  "Deploy pods"),
        (GREEN,  "Wait rollout"),
    ]
    x = Inches(0.55)
    for i, (col, label) in enumerate(flow_items):
        add_text(slide, label, x, Inches(4.93), Inches(1.42), Inches(0.32), size=9.5, color=col, bold=True)
        x += Inches(1.48)
        if i < len(flow_items) - 1:
            add_text(slide, "→", x - Inches(0.12), Inches(4.93), Inches(0.2), Inches(0.32), size=9, color=MUTED)

    add_text(slide,
             "› Secrets injected at runtime — never baked into image or git     "
             "› OIDC: no AWS_ACCESS_KEY_ID stored anywhere     "
             "› Parameterized image tag — controlled versioning",
             Inches(0.55), Inches(5.35), Inches(12.0), Inches(0.55),
             size=9.5, color=TEXT)

    slide_counter(slide, prs, 3, 4)
    return slide


# ── SLIDE 4 ───────────────────────────────────────────────

def make_slide4(prs):
    slide = blank_slide(prs)
    fill_bg(slide, prs)

    slide_label(slide, "Kubernetes — Cluster Internals")
    add_text(slide, "EKS Architecture — Pods, Networking & Permissions",
             Inches(0.42), Inches(0.5), Inches(12), Inches(0.55),
             size=28, bold=True, color=ACCENT)
    add_divider(slide, Inches(1.18))

    col_w = Inches(4.1)
    gap   = Inches(0.22)

    # ── col 1: Networking / routing ──
    x1 = Inches(0.42)
    add_rect(slide, x1, Inches(1.35), col_w, Inches(2.3))
    add_text(slide, "TRAFFIC ROUTING",
             x1 + Inches(0.13), Inches(1.42), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb = slide.shapes.add_textbox(x1 + Inches(0.13), Inches(1.72), col_w - Inches(0.2), Inches(1.85))
    tf = tb.text_frame; tf.word_wrap = True
    add_para(tf, "ALB — internet-facing, HTTP :80", size=10.5, color=TEXT)
    add_para(tf, "Ingress path rules:", size=10.5, color=TEXT)
    p = tf.add_paragraph(); p.font.size = Pt(10); p.font.color.rgb = ACCENT
    p.text = "    /users /pets /appointments … → backend-service :5000"
    p = tf.add_paragraph(); p.font.size = Pt(10); p.font.color.rgb = ACCENT
    p.text = "    / (catch-all)               → frontend-service :80"
    add_para(tf, "ClusterIP services (pods not directly exposed)", size=10.5, color=MUTED)
    add_para(tf, "target-type: ip — routes to pod IPs (AWS VPC CNI)", size=10.5, color=TEXT)

    add_rect(slide, x1, Inches(3.78), col_w, Inches(1.5))
    add_text(slide, "AWS LOAD BALANCER CONTROLLER",
             x1 + Inches(0.13), Inches(3.85), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb2 = slide.shapes.add_textbox(x1 + Inches(0.13), Inches(4.15), col_w - Inches(0.2), Inches(1.05))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    add_para(tf2, "Runs in kube-system namespace", size=10.5, color=TEXT)
    add_para(tf2, "Installed via Helm on every deploy (idempotent)", size=10.5, color=TEXT)
    add_para(tf2, "IRSA: own IAM role scoped to ALB/EC2 APIs only", size=10.5, color=GREEN)

    # ── col 2: Pods & config ──
    x2 = x1 + col_w + gap
    add_rect(slide, x2, Inches(1.35), col_w, Inches(1.65))
    add_text(slide, "BACKEND — 2 REPLICAS  (Flask :5000)",
             x2 + Inches(0.13), Inches(1.42), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb3 = slide.shapes.add_textbox(x2 + Inches(0.13), Inches(1.72), col_w - Inches(0.2), Inches(1.2))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    add_para(tf3, "Image from ECR — tag injected at deploy time", size=10.5, color=TEXT)
    add_para(tf3, "runAsUser=1000, allowPrivilegeEscalation=false", size=10.5, color=GREEN)
    add_para(tf3, "250m/256Mi → 500m/512Mi  |  maxUnavailable=0", size=10.5, color=TEXT)
    add_para(tf3, "Liveness + readiness probes on /health", size=10.5, color=TEXT)

    add_rect(slide, x2, Inches(3.12), col_w, Inches(1.15))
    add_text(slide, "FRONTEND — 2 REPLICAS  (Nginx :80)",
             x2 + Inches(0.13), Inches(3.19), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb4 = slide.shapes.add_textbox(x2 + Inches(0.13), Inches(3.49), col_w - Inches(0.2), Inches(0.7))
    tf4 = tb4.text_frame; tf4.word_wrap = True
    add_para(tf4, "Multi-stage: Node build → Nginx static serve", size=10.5, color=TEXT)
    add_para(tf4, "No secrets — pure static SPA  |  maxUnavailable=0", size=10.5, color=TEXT)

    add_rect(slide, x2, Inches(4.37), col_w, Inches(1.4))
    add_text(slide, "CONFIG & SECRETS",
             x2 + Inches(0.13), Inches(4.44), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb5 = slide.shapes.add_textbox(x2 + Inches(0.13), Inches(4.74), col_w - Inches(0.2), Inches(0.95))
    tf5 = tb5.text_frame; tf5.word_wrap = True
    add_para(tf5, "ConfigMap backend-config — DB_HOST, CORS_ORIGINS (non-sensitive)", size=10.5, color=TEXT)
    add_para(tf5, "Secret backend-secrets — DB_PASSWORD, JWT_SECRET_KEY", size=10.5, color=RED)
    add_para(tf5, "DB migration job: schema.sql + seed_data.sql before pods start", size=10.5, color=TEXT)

    # ── col 3: IAM / permissions ──
    x3 = x2 + col_w + gap
    add_rect(slide, x3, Inches(1.35), col_w, Inches(2.9))
    add_text(slide, "IAM ROLES & PERMISSIONS",
             x3 + Inches(0.13), Inches(1.42), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb6 = slide.shapes.add_textbox(x3 + Inches(0.13), Inches(1.72), col_w - Inches(0.2), Inches(2.45))
    tf6 = tb6.text_frame; tf6.word_wrap = True
    iam_items = [
        ("EKS Node Role",        " → AmazonEC2ContainerRegistryReadOnly\n    (nodes pull ECR — no imagePullSecrets needed)"),
        ("LBC Role (IRSA)",      " → create/manage ALBs, target groups\n    (scoped to ELB + EC2 APIs only)"),
        ("GitHub OIDC Role",     " → AdministratorAccess — CI/CD deploys,\n    pushes ECR, manages Terraform infra"),
        ("EKS Cluster Role",     " → AmazonEKSClusterPolicy"),
    ]
    for key, val in iam_items:
        p = tf6.add_paragraph(); p.font.size = Pt(10.5)
        r1 = p.add_run(); r1.text = "› " + key; r1.font.bold = True; r1.font.color.rgb = PURPLE
        r2 = p.add_run(); r2.text = val; r2.font.color.rgb = TEXT
        tf6.add_paragraph()  # spacer

    add_rect(slide, x3, Inches(4.37), col_w, Inches(1.4))
    add_text(slide, "SECURITY SUMMARY",
             x3 + Inches(0.13), Inches(4.44), col_w - Inches(0.2), Inches(0.28),
             size=8, bold=True, color=MUTED)
    tb7 = slide.shapes.add_textbox(x3 + Inches(0.13), Inches(4.74), col_w - Inches(0.2), Inches(0.95))
    tf7 = tb7.text_frame; tf7.word_wrap = True
    for item in [
        "No imagePullSecrets — node IAM role handles ECR auth",
        "IRSA — LBC gets scoped role, not broad node role",
        "K8s Secrets — sensitive values separate from ConfigMap",
        "RDS unreachable — private subnet + SG from EKS nodes only",
    ]:
        add_para(tf7, item, size=10.5, color=GREEN)

    slide_counter(slide, prs, 4, 4)
    return slide


# ── MAIN ─────────────────────────────────────────────────

def main():
    prs = new_prs()
    make_slide1(prs)
    make_slide2(prs)
    make_slide3(prs)
    make_slide4(prs)
    prs.save(OUT_FILE)
    print(f"Saved: {OUT_FILE}")

if __name__ == "__main__":
    main()
